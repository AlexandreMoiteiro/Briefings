# app.py — Briefings com editor de NOTAMs, GAMET e SIGMET (via Gist)
# + Charts (Weather) + PDFs (Detailed/PT e Final/EN) + PowerPoint (EN)
# + Emparelhamento Navlog↔VFR por Rota + Flight Plan + Mass & Balance
from typing import Dict, Any, List, Tuple, Optional
import io, os, re, base64, tempfile, unicodedata, json, datetime as dt
import streamlit as st
from PIL import Image
from fpdf import FPDF
import fitz  # PyMuPDF
import requests
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt

# ---------- External pages ----------
APP_WEATHER_URL = "https://briefings.streamlit.app/Weather"
APP_NOTAMS_URL  = "https://briefings.streamlit.app/NOTAMs"
APP_VFRMAP_URL  = "https://briefings.streamlit.app/VFRMap"
APP_MNB_URL     = "https://briefings.streamlit.app/MassBalance"
APP_NAV_LOG     = "https://briefings.streamlit.app/NavLog"

# ---------- Página & estilos ----------
st.set_page_config(page_title="Briefings", layout="wide")
st.markdown("""
<style>
:root { --muted:#6b7280; --line:#e5e7eb; --pastel:#5a7fb3; }
.app-top { display:flex; align-items:center; gap:.75rem; flex-wrap:wrap; margin-bottom:.35rem; }
.app-title { font-size: 2.2rem; font-weight: 800; margin: 0 .5rem .25rem 0; }
.small { font-size:.92rem; color:var(--muted); }
.monos{font-family:ui-monospace,Menlo,Consolas,monospace;white-space:pre-wrap}
hr{border:none;border-top:1px solid var(--line);margin:12px 0}
[data-testid="stSidebar"], [data-testid="stSidebarNav"] { display:none !important; }
[data-testid="stSidebarCollapseButton"] { display:none !important; }
header [data-testid="baseButton-headerNoPadding"] { display:none !important; }
.block-label{font-weight:700;margin:.25rem 0}
.section-card{border:1px solid var(--line); border-radius:12px; padding:12px 14px; background:#fff}
.kv{display:grid;grid-template-columns:140px 1fr; gap:.35rem .8rem}
.kv .k{color:#374151}
.kv .v{color:#111827}
.btnbar a{display:inline-block;padding:6px 10px;border:1px solid var(--line);border-radius:8px;text-decoration:none;font-weight:600;color:#111827;background:#f8fafc}
.btnbar a:hover{background:#f1f5f9}
</style>
""", unsafe_allow_html=True)

# ---------- OpenAI client ----------
client = OpenAI(api_key=st.secrets.get("OPENAI_API_KEY"))

# ---------- Constantes úteis ----------
PASTEL = (90, 127, 179)
LPSO_ARP = (39.211667, -8.057778)  # LPSO (ARP)

# ---------- Utils ----------
def safe_str(x) -> str:
    return "" if x is None or x is Ellipsis else str(x)

def ascii_safe(text: str) -> str:
    if text is None or text is Ellipsis:
        return ""
    t = unicodedata.normalize("NFKD", str(text)).encode("ascii", "ignore").decode("ascii")
    return (t.replace("\u00A0"," ").replace("\u2009"," ").replace("\u2013","-")
             .replace("\u2014","-").replace("\uFEFF",""))

def parse_icaos(s: str) -> List[str]:
    tokens = re.split(r"[,\s]+", (s or "").strip(), flags=re.UNICODE)
    return [t.upper() for t in tokens if t]

def read_upload_bytes(upload) -> bytes:
    try:
        return upload.getvalue() if hasattr(upload, "getvalue") else upload.read()
    except Exception:
        return b""

# ---------- Image helpers ----------
def load_first_pdf_page(pdf_bytes: bytes, dpi: int = 450) -> Image.Image:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page = doc.load_page(0)
    png = page.get_pixmap(dpi=dpi).tobytes("png")
    return Image.open(io.BytesIO(png)).convert("RGB").copy()

def ensure_png_from_bytes(file_bytes: bytes, mime: str) -> io.BytesIO:
    """Aceita PDF/PNG/JPG/JPEG/GIF e devolve bytes PNG (ou placeholder)."""
    try:
        m = (mime or "").lower()
        if m == "application/pdf":
            img = load_first_pdf_page(file_bytes, dpi=300)
        else:
            img = Image.open(io.BytesIO(file_bytes)).convert("RGB").copy()
        out = io.BytesIO(); img.save(out, format="PNG"); out.seek(0); return out
    except Exception:
        try:
            Image.open(io.BytesIO(file_bytes))
            return io.BytesIO(file_bytes)
        except Exception:
            ph = Image.new("RGB", (800, 600), (245, 246, 248))
            out = io.BytesIO(); ph.save(out, format="PNG"); out.seek(0); return out

# ---------- Gist helpers ----------
def _get_gist(seckey_token, seckey_id, seckey_file, fallback_file="") -> Tuple[str,str,str]:
    token = (st.secrets.get(seckey_token) or st.secrets.get("GIST_TOKEN") or "").strip()
    gid   = (st.secrets.get(seckey_id)    or st.secrets.get("GIST_ID")    or "").strip()
    fn    = (st.secrets.get(seckey_file)  or fallback_file).strip()
    return token, gid, fn

def _gist_ok(token, gid, fn) -> bool: return all([token, gid, fn])

@st.cache_data(ttl=90)
def gist_load(token, gid, fn) -> Dict[str,Any]:
    if not _gist_ok(token,gid,fn): return {}
    try:
        r = requests.get(f"https://api.github.com/gists/{gid}",
                         headers={"Authorization": f"token {token}", "Accept": "application/vnd.github+json"},
                         timeout=12)
        r.raise_for_status()
        files = r.json().get("files", {})
        file_obj = files.get(fn) or {}
        content = (file_obj.get("content") or "").strip()
        return json.loads(content) if content else {}
    except Exception:
        return {}

def gist_save(token,gid,fn,payload:Dict[str,Any]) -> Tuple[bool,str]:
    if not _gist_ok(token,gid,fn): return False, "Config Gist incompleta."
    try:
        body = {"files": {fn: {"content": json.dumps(payload, ensure_ascii=False, indent=2)}}}
        r = requests.patch(f"https://api.github.com/gists/{gid}",
                           headers={"Authorization": f"token {token}", "Accept": "application/vnd.github+json"},
                           json=body, timeout=12)
        if r.status_code >= 400: return False, f"GitHub respondeu {r.status_code}: {r.text}"
        return True, "Guardado no Gist."
    except Exception as e: return False, f"Erro a gravar no Gist: {e}"

# GAMET
def _get_gamet_secrets(): return _get_gist("GAMET_GIST_TOKEN","GAMET_GIST_ID","GAMET_GIST_FILENAME","gamet.json")
def load_gamet_from_gist() -> Dict[str, Any]:
    t,g,f = _get_gamet_secrets(); d = gist_load(t,g,f)
    return d if d else {"text":"","updated_utc":None}
def save_gamet_to_gist(text: str) -> Tuple[bool, str]:
    t,g,f = _get_gamet_secrets()
    payload = {"updated_utc": dt.datetime.utcnow().strftime("%Y-%m-%dT%H:%MZ"), "text": (text or "").strip()}
    return gist_save(t,g,f,payload)

# NOTAMs
def _get_notam_secrets(): return _get_gist("NOTAM_GIST_TOKEN","NOTAM_GIST_ID","NOTAM_GIST_FILENAME","notams.json")
def load_notams_from_gist() -> Dict[str, Any]:
    t,g,f = _get_notam_secrets(); d = gist_load(t,g,f)
    if not d: return {"map": {}, "updated_utc": None}
    if "map" in d: return {"map": d.get("map") or {}, "updated_utc": d.get("updated_utc")}
    m = {k:v for k,v in d.items() if isinstance(v,list)}; return {"map": m, "updated_utc": d.get("updated_utc")}
def save_notams_to_gist(new_map: Dict[str, List[str]]) -> Tuple[bool, str]:
    t,g,f = _get_notam_secrets()
    payload = {"updated_utc": dt.datetime.utcnow().strftime("%Y-%m-%dT%H:%MZ"),
               "map": {k: [s for s in (v or []) if str(s).strip()] for k,v in new_map.items()}}
    return gist_save(t,g,f,payload)

# SIGMET
def _get_sigmet_secrets(): return _get_gist("SIGMET_GIST_TOKEN","SIGMET_GIST_ID","SIGMET_GIST_FILENAME","sigmet.json")
def load_sigmet_from_gist() -> Dict[str, Any]:
    t,g,f = _get_sigmet_secrets(); d = gist_load(t,g,f)
    return d if d else {"text":"","updated_utc":None}
def save_sigmet_to_gist(text: str) -> Tuple[bool, str]:
    t,g,f = _get_sigmet_secrets()
    payload = {"updated_utc": dt.datetime.utcnow().strftime("%Y-%m-%dT%H:%MZ"), "text": (text or "").strip()}
    return gist_save(t,g,f,payload)

# ---------- METAR/TAF (CheckWX) ----------
def cw_headers() -> Dict[str, str]:
    key = st.secrets.get("CHECKWX_API_KEY", "\n").strip()
    return {"X-API-Key": key} if key else {}

def fetch_metar_now(icao: str) -> str:
    try:
        hdr = cw_headers()
        if not hdr: return ""
        r = requests.get(f"https://api.checkwx.com/metar/{icao}", headers=hdr, timeout=10)
        r.raise_for_status()
        data = r.json().get("data", [])
        if not data: return ""
        if isinstance(data[0], dict): return data[0].get("raw") or data[0].get("raw_text", "") or ""
        return str(data[0])
    except Exception: return ""

def fetch_taf_now(icao: str) -> str:
    try:
        hdr = cw_headers()
        if not hdr: return ""
        r = requests.get(f"https://api.checkwx.com/taf/{icao}", headers=hdr, timeout=10)
        r.raise_for_status()
        data = r.json().get("data", [])
        if not data: return ""
        if isinstance(data[0], dict): return data[0].get("raw") or data[0].get("raw_text", "") or ""
        return str(data[0])
    except Exception: return ""

# ---------- GPT wrapper ----------
def gpt_text(prompt_system: str, prompt_user: str, max_tokens: int = 900) -> str:
    model_name = safe_str(st.secrets.get("OPENAI_MODEL", "gpt-4o-mini")).strip() or "gpt-4o-mini"
    try:
        r2 = client.chat.completions.create(
            model=model_name,
            messages=[{"role":"system","content":prompt_system},{"role":"user","content":prompt_user}],
            max_tokens=max_tokens, temperature=0.15
        )
        content = (r2.choices[0].message.content or "").strip()
        return ascii_safe(content) if content else ""
    except Exception as e2:
        return ascii_safe(f"Analise indisponivel (erro IA: {e2})")

# ---------- Prompts (PT-PT) ----------
def analyze_chart_pt(kind: str, img_b64: str, filename_hint: str = "") -> str:
    model_name = safe_str(st.secrets.get("OPENAI_MODEL_VISION", "gpt-4o")).strip() or "gpt-4o"
    sys = (
        "Es meteorologista aeronautico senior. PT-PT, conciso e rigoroso, texto corrido com 5 blocos curtos. "
        "Usa apenas informacao visivel; se algo nao existir, 'nao indicado'."
    )
    if not (st.secrets.get("OPENAI_API_KEY") or "").strip():
        return "Analise de imagem desativada (OPENAI_API_KEY em falta)."
    try:
        r = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "system", "content": sys},
                {"role": "user", "content": [
                    {"type": "text", "text": f"Tipo de chart: {kind}. Ficheiro: {filename_hint}"},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}}
                ]},
            ],
            max_tokens=1100, temperature=0.1
        )
        out = (r.choices[0].message.content or "").strip()
        return ascii_safe(out) if out else "Analise indisponivel."
    except Exception as e:
        return ascii_safe(f"Analise indisponivel (erro IA: {e})")

def analyze_metar_taf_pt(icao: str, metar: str, taf: str) -> str:
    sys = (
        "Es meteorologista aeronautico senior. PT-PT, resposta telegráfica e concisa (max ~8 linhas). "
        "Usa apenas info presente; sem glossarios. Inclui: hora, vento/raj, vis, fenomenos, nuvens+alturas (oktas entre parenteses), T/Td, QNH. "
        "No TAF: validade e BECMG/TEMPO/PROB com efeito pratico (1 frase/segmento). Se algo nao existir, 'nao presente'. "
        "Termina com 'Impacto' (VFR/IFR + 2-3 riscos)."
    )
    user = f"Aerodromo {icao}\n\nMETAR (RAW):\n{metar}\n\nTAF (RAW):\n{taf}"
    return gpt_text(prompt_system=sys, prompt_user=user, max_tokens=700)

# ---------- PDF helpers ----------
def draw_header(pdf: FPDF, text: str) -> None:
    pdf.set_draw_color(229, 231, 235)
    pdf.set_line_width(0.3)
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, ascii_safe(text), ln=True, align="C", border="B")

def place_image_full(pdf: FPDF, img_png: io.BytesIO, max_h_pad: int = 58) -> None:
    max_w = pdf.w - 22; max_h = pdf.h - max_h_pad
    img = Image.open(img_png); iw, ih = img.size
    r = min(max_w / iw, max_h / ih); w, h = int(iw * r), int(ih * r)
    x = (pdf.w - w) // 2; y = pdf.get_y() + 6
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        img.save(tmp, format="PNG"); path = tmp.name
    pdf.image(path, x=x, y=y, w=w, h=h); os.remove(path)
    pdf.ln(h + 10)

def image_bytes_to_pdf_bytes_fullbleed(img_bytes: bytes, orientation: str = "P") -> bytes:
    """Imagem -> 1 página PDF full-bleed (sem títulos/margens)."""
    doc = FPDF(orientation=orientation, unit="mm", format="A4"); doc.add_page(orientation=orientation)
    max_w, max_h = doc.w, doc.h
    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    iw, ih = img.size; r = min(max_w/iw, max_h/ih); w, h = int(iw*r), int(ih*r)
    x, y = (doc.w - w) / 2, (doc.h - h) / 2
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        img.save(tmp, "PNG"); path = tmp.name
    doc.image(path, x=x, y=y, w=w, h=h); os.remove(path)
    data = doc.output(dest="S")
    return data if isinstance(data, (bytes, bytearray)) else str(data).encode("latin-1")

def fpdf_to_bytes(doc: FPDF) -> bytes:
    data = doc.output(dest="S")
    return data if isinstance(data, (bytes, bytearray)) else str(data).encode("latin-1")

# ---------- Ordenação de charts ----------
_KIND_RANK = {"SPC": 1, "SIGWX": 2, "Wind & Temp": 3, "Other": 9}
def _chart_sort_key(c: Dict[str, Any]) -> Tuple[int, int]:
    kind = c.get("kind", "Other")
    rank = _KIND_RANK.get(kind, 9)
    order = int(c.get("order", 9999) or 9999)
    return (rank, order)

# ---------- PDF classes ----------
class DetailedPDF(FPDF):
    def header(self) -> None: pass
    def footer(self) -> None: pass

    def chart_block(self, title: str, subtitle: str, img_png: io.BytesIO, analysis_pt: str) -> None:
        self.add_page(orientation="P"); draw_header(self, ascii_safe(title))
        if subtitle: self.set_font("Helvetica","I",12); self.cell(0,9,ascii_safe(subtitle), ln=True, align="C")
        max_w = self.w - 22; max_h = (self.h // 2) - 18
        img = Image.open(img_png); iw, ih = img.size; r = min(max_w/iw, max_h/ih); w, h = int(iw*r), int(ih*r)
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            img.save(tmp, format="PNG"); path = tmp.name
        self.image(path, x=(self.w - w)//2, y=self.get_y() + 6, w=w, h=h); os.remove(path); self.ln(h + 12)
        self.set_font("Helvetica","",12); self.multi_cell(0, 7, ascii_safe(analysis_pt or " "))

class FinalBriefPDF(FPDF):
    def header(self) -> None: pass
    def footer(self) -> None: pass
    def cover(self, mission_no, pilot, aircraft, callsign, reg, date_str, time_utc) -> None:
        self.add_page(orientation="L"); self.set_xy(0, 36)
        self.set_font("Helvetica","B",28); self.cell(0,14,"Briefing", ln=True, align="C")
        self.ln(2); self.set_font("Helvetica","",13)
        self.cell(0,8,ascii_safe(f"Mission: {mission_no}"), ln=True, align="C")
        if pilot or aircraft or callsign or reg:
            self.cell(0,8,ascii_safe(f"Pilot: {pilot}   Aircraft: {aircraft}   Callsign: {callsign}   Reg: {reg}"), ln=True, align="C")
        if date_str or time_utc:
            self.cell(0,8,ascii_safe(f"Date: {date_str}   UTC: {time_utc}"), ln=True, align="C")
        self.ln(4); self.set_text_color(*PASTEL); self.set_font("Helvetica","",12)
        self.cell(0,8, ascii_safe("Weather page: ") + APP_WEATHER_URL, ln=True, align="C", link=APP_WEATHER_URL)
        self.cell(0,8, ascii_safe("NOTAMs page: ") + APP_NOTAMS_URL, ln=True, align="C", link=APP_NOTAMS_URL)
        self.set_text_color(0,0,0)
    def charts_only(self, charts: List[Tuple[str, str, io.BytesIO]]) -> None:
        for (title, subtitle, img_png) in charts:
            self.add_page(orientation="L"); draw_header(self, ascii_safe(title))
            if subtitle:
                self.set_font("Helvetica","I",12); self.cell(0,9,ascii_safe(subtitle), ln=True, align="C")
            place_image_full(self, img_png)

# ---------- UI topo ----------
st.markdown(
    f'''<div class="app-top"><div class="app-title">Briefings</div>
    <span class="btnbar">
      <a href="{APP_WEATHER_URL}" target="_blank">Weather</a>
      <a href="{APP_NOTAMS_URL}" target="_blank">NOTAMs</a>
      <a href="{APP_VFRMAP_URL}" target="_blank">VFR Map</a>
      <a href="{APP_MNB_URL}" target="_blank">Mass & Balance</a>
      <a href="{APP_NAV_LOG}" target="_blank">NavLog </a>
    </span></div>''', unsafe_allow_html=True
)

# ---------- Abas ----------
tab_mission, tab_notams, tab_sigmet_gamet, tab_charts, tab_pairs, tab_fpmb, tab_generate = st.tabs(
    ["Missão", "NOTAMs", "SIGMET & GAMET", "Charts", "Navlog ↔ VFR (Rotas)", "Flight Plan & M&B", "Gerar Saídas"]
)

# ---------- Missão ----------
with tab_mission:
    st.markdown("### Dados da Missão")
    colA, colB, colC = st.columns(3)
    with colA:
        pilot = st.text_input("Pilot name", "Alexandre Moiteiro")
        callsign = st.text_input("Mission callsign", "")
    with colB:
        aircraft_type = st.text_input("Aircraft type", "Tecnam P2008")
        registration = st.text_input("Registration", "CS-XXX")
    with colC:
        mission_no = st.text_input("Mission number", "")
        flight_date = st.date_input("Flight date")
        time_utc = st.text_input("UTC time", "")

    st.markdown("#### Aerodromes")
    c1, c2 = st.columns(2)
    with c1:
        icaos_metar_str = st.text_input("ICAO list for METAR/TAF (comma / space / newline)", value="LPPT LPBJ LEBZ")
        icaos_metar = parse_icaos(icaos_metar_str)
    with c2:
        icaos_notam_str = st.text_input("ICAO list for NOTAMs (comma / space / newline)", value="LPPC(Enroute) LPSO LPCB LPEV")
        icaos_notam = parse_icaos(icaos_notam_str)

# ---------- NOTAMs ----------
with tab_notams:
    st.markdown("### NOTAMs (editar e guardar)")
    saved_notams = load_notams_from_gist()
    existing_map: Dict[str, List[str]] = (saved_notams.get("map") or {}) if isinstance(saved_notams, dict) else {}

    def parse_block_to_list(text: str) -> List[str]:
        if not (text or "").strip(): return []
        parts = re.split(r"\n\s*\n+", text.strip())
        return [p.strip() for p in parts if p.strip()]

    edit_cols = st.columns(3)
    editors_notam: Dict[str, str] = {}
    for i, icao in enumerate(icaos_notam if 'icaos_notam' in locals() else []):
        with edit_cols[i % 3]:
            initial_text = "\n\n".join(existing_map.get(icao, [])) if existing_map.get(icao) else ""
            editors_notam[icao] = st.text_area(
                f"{icao} — NOTAMs",
                value=initial_text,
                placeholder=("Ex.: AERODROME BEACON ONLY FLASH-GREEN LIGHT OPERATIVE.\n"
                             "FROM: 29th Jul 2025 15:10 TO: 29th Sep 2025 18:18 EST\n\n"
                             "Outro NOTAM aqui..."),
                key=f"ed_notam_{i}",
                height=160
            )

    if st.button("Guardar NOTAMs no Gist"):
        new_map: Dict[str, List[str]] = {}; new_map.update(existing_map)  # merge
        for icao in (icaos_notam if 'icaos_notam' in locals() else []):
            new_map[icao] = parse_block_to_list(editors_notam.get(icao, ""))
        ok, msg = save_notams_to_gist(new_map)
        st.success(msg) if ok else st.error(msg)
        try: st.cache_data.clear()
        except Exception: pass

# ---------- SIGMET & GAMET ----------
with tab_sigmet_gamet:
    st.markdown("### SIGMET & GAMET")
    _gamet_obj = load_gamet_from_gist(); _gamet_initial = safe_str((_gamet_obj.get("text") or "").strip())
    gamet_text = st.text_area("GAMET — texto integral", value=_gamet_initial, height=220, key="gamet_editor",
                              placeholder="Ex.: LPPC FIR GAMET VALID 12/06Z-12/12Z\n...")
    if st.button("Guardar GAMET no Gist"):
        ok, msg = save_gamet_to_gist(gamet_text); st.success(msg) if ok else st.error(msg)
        try: st.cache_data.clear()
        except Exception: pass

    st.divider()

    _sigmet_obj = load_sigmet_from_gist(); _sigmet_initial = safe_str((_sigmet_obj.get("text") or "").strip())
    sigmet_text = st.text_area("SIGMET (LPPC) — texto integral", value=_sigmet_initial, height=160, key="sigmet_editor",
                               placeholder="Ex.: LPPC SIGMET 2 VALID 12/09Z-12/13Z LPPC-\nSEV TURB FCST BTN FL080/FL200 MOV NE 20KT ...")
    if st.button("Guardar SIGMET no Gist"):
        ok, msg = save_sigmet_to_gist(sigmet_text); st.success(msg) if ok else st.error(msg)
        try: st.cache_data.clear()
        except Exception: pass

# ---------- Charts ----------
with tab_charts:
    st.markdown("### Charts (SIGWX / SPC / Wind & Temp / Other)")
    st.caption("Aceita PDF/PNG/JPG/JPEG/GIF (para PDF lemos a 1.ª página).")
    use_ai_for_charts = st.toggle("Analisar charts com IA", value=True, help="Marcado por omissao")
    preview_w = st.slider("Largura da pré-visualização (px)", min_value=240, max_value=640, value=420, step=10)
    uploads = st.file_uploader("Upload charts", type=["pdf","png","jpg","jpeg","gif"], accept_multiple_files=True, label_visibility="collapsed")

    def _base_title_for_kind(k: str) -> str:
        return {"SIGWX":"Significant Weather Chart (SIGWX)","SPC":"Surface Pressure Chart (SPC)","Wind & Temp":"Wind and Temperature Chart"}.get(k, "Weather Chart")

    charts: List[Dict[str,Any]] = []
    if uploads:
        for idx, f in enumerate(uploads):
            raw = f.read(); mime = f.type or ""
            img_png = ensure_png_from_bytes(raw, mime)
            name = safe_str(getattr(f, "name", "")) or "(sem nome)"
            col_img, col_meta = st.columns([0.35, 0.65])
            with col_img:
                try: st.image(img_png.getvalue(), caption=name, width=preview_w)
                except Exception: st.write(name)
            with col_meta:
                kind = st.selectbox(f"Tipo do chart #{idx+1}", ["SIGWX","SPC","Wind & Temp","Other"], index=0, key=f"kind_{idx}")
                title_default = _base_title_for_kind(kind)
                title = st.text_input("Título", value=title_default, key=f"title_{idx}")
                subtitle = st.text_input("Subtítulo (opcional)", value="", key=f"subtitle_{idx}")
                order_val = st.number_input("Ordem", min_value=1, max_value=len(uploads)+10, value=idx+1, step=1, key=f"ord_{idx}")
            charts.append({"kind": kind, "title": title, "subtitle": subtitle, "img_png": img_png, "order": order_val, "filename": name})

# ---------- Navlog ↔ VFR (pares por ROTA) ----------
with tab_pairs:
    st.markdown("### Emparelhamento Navlog ↔ VFR por ROTA")
    st.caption("Para cada ROTA (ex.: LPSO-LPCB, LPSO-LPEV) carrega um Navlog e o respetivo mapa VFR. Aceita PDF/PNG/JPG/JPEG/GIF.")
    num_pairs = st.number_input("Número de pares (Rotas)", min_value=0, max_value=8, value=0, step=1)
    pairs: List[Dict[str, Any]] = []
    for i in range(int(num_pairs)):
        with st.expander(f"Rota #{i+1}", expanded=False):
            route = safe_str(st.text_input("ROTA (ex.: LPSO-LPCB)", key=f"pair_route_{i}")).upper().strip()
            c1, c2 = st.columns(2)
            with c1:
                nav_file = st.file_uploader(f"Navlog ({route or 'ROTA'})", type=["pdf","png","jpg","jpeg","gif"], key=f"pair_nav_{i}")
            with c2:
                vfr_file = st.file_uploader(f"VFR Map ({route or 'ROTA'})", type=["pdf","png","jpg","jpeg","gif"], key=f"pair_vfr_{i}")
            pairs.append({"route": route, "nav": nav_file, "vfr": vfr_file})

# ---------- Flight Plan & M&B ----------
with tab_fpmb:
    st.markdown("### Flight Plan & Mass & Balance")
    c1, c2 = st.columns(2)
    with c1:
        fp_upload = st.file_uploader("Flight Plan (PDF/PNG/JPG)", type=["pdf","png","jpg","jpeg"])
        if fp_upload: st.success(f"Flight Plan carregado: {safe_str(fp_upload.name)}")
    with c2:
        mb_upload = st.file_uploader("Mass & Balance (PDF/PNG/JPG)", type=["pdf","png","jpg","jpeg"])
        if mb_upload: st.success(f"M&B carregado: {safe_str(mb_upload.name)}")

# ---------- Gerar Saídas ----------
with tab_generate:
    st.markdown("### Saídas")
    col_pdfs = st.columns(3)
    with col_pdfs[0]: gen_det   = st.button("Generate Detailed (PT)")
    with col_pdfs[1]: gen_final = st.button("Generate Final Briefing (EN)")
    with col_pdfs[2]: gen_ppt   = st.button("Generate PowerPoint (EN)")

# ---------- Detailed (PT) ----------
if 'gen_det' in locals() and gen_det:
    det = DetailedPDF()

    # CHARTS PRIMEIRO (Weather=charts)
    charts_local: List[Dict[str,Any]] = locals().get("charts", [])
    if charts_local:
        grouped: Dict[str, List[Dict[str,Any]]] = {"SPC": [], "SIGWX": [], "Wind & Temp": [], "Other": []}
        for c in charts_local: grouped.setdefault(c["kind"], []).append(c)
        for k in list(grouped.keys()): grouped[k] = sorted(grouped[k], key=_chart_sort_key)
        for kind in ["SPC","SIGWX","Wind & Temp","Other"]:
            for ch in grouped.get(kind, []):
                title, subtitle, img_png, fname = ch["title"], ch["subtitle"], ch["img_png"], ch.get("filename","")
                analysis_txt = ""
                if locals().get("use_ai_for_charts", False):
                    try:
                        analysis_txt = analyze_chart_pt(kind, base64.b64encode(img_png.getvalue()).decode("utf-8"), filename_hint=fname)
                    except Exception: analysis_txt = "Analise indisponivel."
                det.chart_block(title, subtitle, img_png, analysis_txt)

    # (Opcional) METAR/TAF resumido — se quiseres manter no Detailed
    icaos_metar_local = locals().get("icaos_metar", [])
    if icaos_metar_local:
        det.add_page(orientation="P"); draw_header(det, "METAR / TAF — Interpretacao (PT, resumida)")
        det.set_font("Helvetica","",12); det.ln(2)
        for icao in icaos_metar_local:
            metar_raw = fetch_metar_now(icao) or ""; taf_raw = fetch_taf_now(icao) or ""
            analysis = analyze_metar_taf_pt(icao, metar_raw, taf_raw) if (metar_raw or taf_raw) else "Sem METAR/TAF disponiveis."
            det.set_font("Helvetica","B",13); det.cell(0,8, ascii_safe(icao), ln=True)
            if metar_raw: det.set_font("Helvetica","B",12); det.cell(0,7,"METAR (RAW):", ln=True); det.set_font("Helvetica","",12); det.multi_cell(0,7, ascii_safe(metar_raw))
            if taf_raw:   det.set_font("Helvetica","B",12); det.cell(0,7,"TAF (RAW):", ln=True);   det.set_font("Helvetica","",12); det.multi_cell(0,7, ascii_safe(taf_raw))
            det.set_font("Helvetica","B",12); det.cell(0,7,"Interpretacao:", ln=True)
            det.set_font("Helvetica","",12); det.multi_cell(0,7, ascii_safe(analysis)); det.ln(2)

    # SIGMET / GAMET (texto simples)
    sigmet_text_local = locals().get("sigmet_text",""); _sigmet_initial_local = locals().get("_sigmet_initial","")
    gamet_text_local  = locals().get("gamet_text","");  _gamet_initial_local  = locals().get("_gamet_initial","")
    sigmet_for_pdf = (sigmet_text_local or _sigmet_initial_local or "").strip()
    gamet_for_pdf  = (gamet_text_local  or _gamet_initial_local  or "").strip()
    if sigmet_for_pdf:
        det.add_page(orientation="P"); draw_header(det, "SIGMET (LPPC) — Texto (RAW)")
        det.set_font("Helvetica","",12); det.multi_cell(0,7, ascii_safe(sigmet_for_pdf))
    if gamet_for_pdf:
        det.add_page(orientation="P"); draw_header(det, "GAMET — Texto (RAW)")
        det.set_font("Helvetica","",12); det.multi_cell(0,7, ascii_safe(gamet_for_pdf))

    # Glossário curto
    det.add_page(orientation="P"); draw_header(det, "Glossario — Simbologia, Nuvens & Fenomenos")
    det.set_font("Helvetica","",12); det.multi_cell(0,7, ascii_safe(
        "Cobertura (oktas): FEW 1–2; SCT 3–4; BKN 5–7; OVC 8.\n"
        "Frentes: fria/ quente/ oclusao/ estacionaria. SIGWX: jatos, turbulencia, gelo, CB/TCU, EMBD/OCNL/FRQ."
    ))

    det_name = f"Briefing Detalhado - Missao {locals().get('mission_no') or 'X'}.pdf"
    det.output(det_name)
    with open(det_name, "rb") as f:
        st.download_button("Download Detailed (PT)", data=f.read(), file_name=det_name, mime="application/pdf", use_container_width=True)

# ---------- Final Briefing (EN) ----------
if 'gen_final' in locals() and gen_final:
    fb = FinalBriefPDF()
    fb.cover(
        mission_no=locals().get("mission_no",""),
        pilot=locals().get("pilot",""),
        aircraft=locals().get("aircraft_type",""),
        callsign=locals().get("callsign",""),
        reg=locals().get("registration",""),
        date_str=str(locals().get("flight_date","")),
        time_utc=locals().get("time_utc","")
    )
    # CHARTS primeiro (Weather=charts)
    charts_local: List[Dict[str,Any]] = locals().get("charts", [])
    if charts_local:
        ordered = [(c["title"], c["subtitle"], c["img_png"]) for c in sorted(charts_local, key=_chart_sort_key)]
        fb.charts_only(ordered)

    # Base PDF (forçar bytes)
    fb_bytes: bytes = fpdf_to_bytes(fb)
    fb_bytes = bytes(fb_bytes)  # garante tipo bytes
    final_bytes = fb_bytes

    # Merge com Flight Plan, Pares Navlog/VFR e M&B (tal como são)
    nav_pairs: List[Dict[str, Any]] = locals().get("pairs", [])
    fp_upload = locals().get("fp_upload", None)
    mb_upload = locals().get("mb_upload", None)
    try:
        main = fitz.open(stream=fb_bytes, filetype="pdf"); insert_pos = main.page_count

        # Flight Plan
        if fp_upload is not None:
            raw = read_upload_bytes(fp_upload)
            if (fp_upload.type or "").lower() == "application/pdf":
                fp_doc = fitz.open(stream=raw, filetype="pdf")
                main.insert_pdf(fp_doc, start_at=insert_pos); insert_pos += fp_doc.page_count; fp_doc.close()
            else:
                fp_bytes = image_bytes_to_pdf_bytes_fullbleed(raw or b"", orientation="P")
                fp_doc = fitz.open(stream=fp_bytes, filetype="pdf")
                main.insert_pdf(fp_doc, start_at=insert_pos); insert_pos += fp_doc.page_count; fp_doc.close()

        # Pares Navlog/VFR
        for p in (nav_pairs or []):
            nv = p.get("nav"); vf = p.get("vfr")
            if nv is not None:
                raw = read_upload_bytes(nv)
                if (nv.type or "").lower() == "application/pdf":
                    nv_doc = fitz.open(stream=raw, filetype="pdf")
                else:
                    nv_doc = fitz.open(stream=image_bytes_to_pdf_bytes_fullbleed(raw or b"", "P"), filetype="pdf")
                main.insert_pdf(nv_doc, start_at=insert_pos); insert_pos += nv_doc.page_count; nv_doc.close()
            if vf is not None:
                raw = read_upload_bytes(vf)
                if (vf.type or "").lower() == "application/pdf":
                    vf_doc = fitz.open(stream=raw, filetype="pdf")
                else:
                    vf_doc = fitz.open(stream=image_bytes_to_pdf_bytes_fullbleed(raw or b"", "L"), filetype="pdf")
                main.insert_pdf(vf_doc, start_at=insert_pos); insert_pos += vf_doc.page_count; vf_doc.close()

        # Mass & Balance
        if mb_upload is not None:
            raw = read_upload_bytes(mb_upload)
            if (mb_upload.type or "").lower() == "application/pdf":
                mb_doc = fitz.open(stream=raw, filetype="pdf")
                main.insert_pdf(mb_doc, start_at=insert_pos); insert_pos += mb_doc.page_count; mb_doc.close()
            else:
                mb_bytes = image_bytes_to_pdf_bytes_fullbleed(raw or b"", orientation="P")
                mb_doc = fitz.open(stream=mb_bytes, filetype="pdf")
                main.insert_pdf(mb_doc, start_at=insert_pos); insert_pos += mb_doc.page_count; mb_doc.close()

        final_bytes = main.tobytes()
        final_bytes = bytes(final_bytes)  # força bytes, mesmo que memoryview/bytearray
        main.close()
    except Exception:
        # mesmo que falhe o merge, garantir bytes na saída base
        final_bytes = bytes(final_bytes)

    final_bytes = bytes(final_bytes)  # redundância segura
    final_name = f"Briefing - Missao {locals().get('mission_no') or 'X'}.pdf"
    st.download_button("Download Final Briefing (EN)", data=final_bytes, file_name=final_name, mime="application/pdf", use_container_width=True)

# ---------- PowerPoint (EN) ----------
def _ppt_blank_slide(prs: Presentation):
    return prs.slide_layouts[6]  # blank

def _ppt_add_fullbleed_image_slide(prs: Presentation, img_bytes: bytes) -> None:
    slide = prs.slides.add_slide(_ppt_blank_slide(prs))
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        Image.open(io.BytesIO(img_bytes)).convert("RGB").save(tmp, "PNG"); path = tmp.name
    # inserir por largura; se exceder a altura, ajusta
    pic = slide.shapes.add_picture(path, Inches(0), Inches(0), width=prs.slide_width)
    if pic.height > prs.slide_height:
        ratio = prs.slide_height / pic.height
        pic.height = prs.slide_height
        pic.width  = int(pic.width * ratio)
        pic.left   = int((prs.slide_width - pic.width) / 2)
    os.remove(path)

def _render_pdf_to_png_bytes_list(pdf_bytes: bytes, dpi: int = 300) -> List[bytes]:
    out = []; doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    for i in range(doc.page_count):
        p = doc.load_page(i); out.append(p.get_pixmap(dpi=dpi).tobytes("png"))
    doc.close(); return out

if 'gen_ppt' in locals() and gen_ppt:
    prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # Capa
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    tf = slide.shapes.title.text_frame; tf.clear(); p = tf.paragraphs[0]
    p.text = "Briefing"; p.font.size = Pt(44); p.font.bold = True
    box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(2)).text_frame
    box.text = ascii_safe(
        f"Mission: {safe_str(locals().get('mission_no',''))}   "
        f"Pilot: {safe_str(locals().get('pilot',''))}   "
        f"Aircraft: {safe_str(locals().get('aircraft_type',''))}   "
        f"Callsign: {safe_str(locals().get('callsign',''))}   "
        f"Reg: {safe_str(locals().get('registration',''))}\n"
        f"Date: {safe_str(str(locals().get('flight_date','')))}   UTC: {safe_str(locals().get('time_utc',''))}\n"
        f"Weather: {APP_WEATHER_URL}    NOTAMs: {APP_NOTAMS_URL}"
    )

    # Charts (Weather)
    charts_local: List[Dict[str,Any]] = locals().get("charts", [])
    if charts_local:
        for c in sorted(charts_local, key=_chart_sort_key):
            _ppt_add_fullbleed_image_slide(prs, c["img_png"].getvalue())

    # Flight Plan
    fp_upload = locals().get("fp_upload", None)
    if fp_upload is not None:
        raw = read_upload_bytes(fp_upload)
        if (fp_upload.type or "").lower() == "application/pdf":
            for page_png in _render_pdf_to_png_bytes_list(raw, dpi=300):
                _ppt_add_fullbleed_image_slide(prs, page_png)
        else:
            _ppt_add_fullbleed_image_slide(prs, raw)

    # Navlog & VFR
    nav_pairs: List[Dict[str, Any]] = locals().get("pairs", [])
    for p in (nav_pairs or []):
        nv = p.get("nav"); vf = p.get("vfr")
        if nv is not None:
            raw = read_upload_bytes(nv)
            if (nv.type or "").lower() == "application/pdf":
                for page_png in _render_pdf_to_png_bytes_list(raw, dpi=300):
                    _ppt_add_fullbleed_image_slide(prs, page_png)
            else:
                _ppt_add_fullbleed_image_slide(prs, raw)
        if vf is not None:
            raw = read_upload_bytes(vf)
            if (vf.type or "").lower() == "application/pdf":
                for page_png in _render_pdf_to_png_bytes_list(raw, dpi=300):
                    _ppt_add_fullbleed_image_slide(prs, page_png)
            else:
                _ppt_add_fullbleed_image_slide(prs, raw)

    # Mass & Balance
    mb_upload = locals().get("mb_upload", None)
    if mb_upload is not None:
        raw = read_upload_bytes(mb_upload)
        if (mb_upload.type or "").lower() == "application/pdf":
            for page_png in _render_pdf_to_png_bytes_list(raw, dpi=300):
                _ppt_add_fullbleed_image_slide(prs, page_png)
        else:
            _ppt_add_fullbleed_image_slide(prs, raw)

    # Guardar PPTX (forçar bytes)
    ppt_name = f"Briefing - Mission {safe_str(locals().get('mission_no') or 'X')}.pptx"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name); tmp.seek(0); ppt_bytes = tmp.read(); tmp_path = tmp.name
    ppt_bytes = bytes(ppt_bytes)  # garante tipo bytes
    st.download_button("Download PowerPoint (EN)", data=ppt_bytes, file_name=ppt_name,
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                       use_container_width=True)
    try: os.remove(tmp_path)
    except Exception: pass
